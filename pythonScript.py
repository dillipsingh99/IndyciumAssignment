"""
Code Description:
used packages : opencv2, numpy, os, python-pptx.
first - I have opened image file with os, after that read the image 
        file widht and height with opencv2, after that i just created logo
        the same width and hieght as image has with numpy, after that i 
        just merge both logo and image with logo position to top left with
        opencv2 addWeighted function.
Second - But while merging the logo and image, with default logo size it almost 
        fit to 1/4th part of the image, so i defined a function named scale to
        reduce the size of logo to fit as expected also i have create one image
        watermark with logo with its default size, sample is within file lol folder.
Third  - Using python-pptx i have created ppt using slide_laout(6) and for loop.
"""

import os
import cv2 as cv
import numpy as np
from pptx import Presentation 
from pptx.util import Inches, Pt

# path of original image
DirPath = '..\pythonAssignment'

# creating list of all original image
files = os.listdir(DirPath)

# filtering original image with extension .jpg using filter and lambda function
image_files = list(filter(lambda x : '.jpg' in x, files))


# function to reduce size of logo to fit over the original image
def scale(image, scale_width):
    (image_height, image_width) = image.shape[:2]
    new_height = int(scale_width / image_width *image_height)
    return cv.resize(image, (scale_width, new_height))

# iterating through image list i.e. files
for file in image_files:
    # Joining files path with DirPath
    imgPath = os.path.join(DirPath, file)
    # Reading image with opencv
    image = cv.imread(imgPath)
    # extracting height and width of image with shape method
    (h_image, w_image) = image.shape[:2]


    """
    The cvtColor() function in OpenCV takes two parameters namely image 
    and code where the image is the image whose color space is to be 
    converted to different color space and the code represents the color conversion code.
    """
    image = cv.cvtColor(image, cv.COLOR_BGR2BGRA)
    # image = cv.cvtColor(image, cv.COLOR_BGR2HSV)

    logoPath = "nike_black.png"
    logo = scale(cv.imread(logoPath, cv.IMREAD_UNCHANGED),1400)
    h_logo, w_logo= logo.shape[:2]
   
    # using numpy overlaying logo on image
    overlay = np.zeros((h_image, w_image, 4), dtype='uint8')
    overlay[0:h_logo, 0:w_logo] = logo
    (o_height, o_width) = overlay.shape[:2]
 
   
    
    output2 = image.copy()
    """
    Following is the syntax of addWeighted() function.

    dst = cv.addWeighted(src1, alpha, src2, beta, gamma[, dst[, dtype]])
    where src1 and src2 are input image arrays and alpha, beta are the 
    corresponding weights to be considered while performing the weighted addition. 
    gamma is static weight that will be added to all the pixels of the image.

    addWeighted() function returns numpy array containing pixel values of the 
    resulting image.
    """
    
    cv.addWeighted(overlay, -25, image, 1.0, 0, output2)
    # writing files to the current directory
    cv.imwrite(f"final{file}", output2)
    
    print(f"{file} processing done.")
    cv.waitKey(0)

# Initializing ppt presentation
prs = Presentation() 
# Creating slide layout as blank page
blank_slide_layout = prs.slide_layouts[6] 
# iterating through range 1-6
for i in range(1,6):
    # Creating slide objects to add ppt
    slide = prs.slides.add_slide(blank_slide_layout)
    # Takinng margin from top, left of slide
    left=Inches(2)
    top=Inches(0.1)
    # textBox for Title
    # adding textBox with add_textbox function takes four argument
    txBox = slide.shapes.add_textbox(left, top, width=Inches(5), height=Inches(1))
    tf = txBox.text_frame
    p = tf.add_paragraph()
    p.font.bold = True
    p.font.size = Pt(40)
    p.text=f"Sample Title {i}"
    # TextBox for Subtile
    # Again adding textBox with add_textbox function takes four argument
    top = Inches(1)
    left = Inches(1)
    stBox = slide.shapes.add_textbox(left, top, width=Inches(5), height=Inches(1))
    st = stBox.text_frame
    s = st.add_paragraph()
    s.font.size = Pt(30)
    s.font.bullet = True
    s.text = f"Sample Subtitle {i}"
    # Adding image to ppt using add_picture function
    # takes same four argument top, left, by default width and height set to zero.
    left = Inches(1)
    top = Inches(2)
    pic = slide.shapes.add_picture(f"finalimage{i}.jpg", left, top, height = Inches(3.5), width=Inches(6)) 
    # incrementing value of i with 1.
    i += 1
# saving presentation
prs.save('pythonProject.pptx')
print("pythonProject.ppt saved to current directory.")